#!/usr/bin/env python3
"""Generate report figures, heuristic comparisons, and Office deliverables."""

from __future__ import annotations

import json
import shutil
import ssl
import sys
import urllib.request
from pathlib import Path

import geopandas as gpd
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.shared import Inches as DocxInches
from docx.text.paragraph import Paragraph
from matplotlib import cm, colors
from matplotlib.lines import Line2D
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt
from scipy import sparse

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from src.evaluation.alternative_design import (  # noqa: E402
    coverage_fraction,
    geographic_spread_selection,
    greedy_max_arcs_selection,
    load_candidate_coords,
)


DATA_PROCESSED = PROJECT_ROOT / "data" / "processed"
FIGURES_DIR = DATA_PROCESSED / "figures"
ARTIFACT_CACHE = Path.home() / ".cache" / "project_argus"
METADATA_200 = DATA_PROCESSED / "visibility_metadata_200.json"
UNIFORM_RESULT = DATA_PROCESSED / "optimization_result_200.json"
RASTER_RESULT = DATA_PROCESSED / "optimization_result_population_raster.json"
REPORT_SOURCE = PROJECT_ROOT / "argus_final_report_v2_edited.docx"
REPORT_OUTPUT = PROJECT_ROOT.parent / "argus_final_report.docx"
PRESENTATION_PATH = PROJECT_ROOT.parent / "argus_presentation.pptx"
PRESENTATION_BACKUP = ARTIFACT_CACHE / "argus_presentation.original.pptx"


def _load_visibility_matrix() -> sparse.csr_matrix:
    for path in (
        DATA_PROCESSED / "visibility_200.npz",
        DATA_PROCESSED / "sensitivity" / "elev_25" / "visibility.npz",
        DATA_PROCESSED / "visibility.npz",
    ):
        if path.exists():
            matrix = sparse.load_npz(path).tocsr()
            if matrix.shape[1] == 200:
                return matrix
    raise FileNotFoundError("Unable to locate a 200-site 25-degree visibility matrix")


def _load_world_basemap() -> gpd.GeoDataFrame:
    cache_dir = ARTIFACT_CACHE / "basemap"
    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_path = cache_dir / "ne_110m_admin_0_countries.zip"
    if not cache_path.exists():
        url = "https://naturalearth.s3.amazonaws.com/110m_cultural/ne_110m_admin_0_countries.zip"
        context = ssl._create_unverified_context()
        with urllib.request.urlopen(url, context=context) as response:
            cache_path.write_bytes(response.read())
    return gpd.read_file(f"zip://{cache_path}")


def _heatmap_dataframe(csv_path: Path) -> tuple[np.ndarray, np.ndarray]:
    frame = pd.read_csv(csv_path)
    elevations = [0.0, 10.0, 25.0]
    budgets = [5, 10, 15, 20, 30]
    coverage = np.zeros((len(elevations), len(budgets)), dtype=float)
    visibility_ub = np.zeros(len(elevations), dtype=float)

    for row_idx, elevation in enumerate(elevations):
        elevation_frame = frame[frame["elevation_deg"].astype(float) == elevation]
        if elevation_frame.empty:
            raise ValueError(f"Missing sensitivity row for elevation {elevation}")
        visibility_ub[row_idx] = 100.0 * float(elevation_frame.iloc[0]["demand_visibility_upper_bound"])
        for col_idx, budget in enumerate(budgets):
            row = elevation_frame[elevation_frame["max_ground_stations"].astype(int) == budget]
            if row.empty:
                raise ValueError(f"Missing sensitivity row for elevation {elevation}, budget {budget}")
            coverage[row_idx, col_idx] = 100.0 * float(row.iloc[0]["achieved_coverage"])
    return coverage, visibility_ub


def build_sensitivity_heatmap(csv_path: Path, output_path: Path, title: str) -> None:
    coverage, visibility_ub = _heatmap_dataframe(csv_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    fig, ax = plt.subplots(figsize=(8, 5), dpi=150)
    coverage_norm = colors.Normalize(vmin=0.0, vmax=100.0)
    ub_norm = colors.Normalize(vmin=0.0, vmax=100.0)
    coverage_cmap = plt.colormaps["Blues"]
    ub_cmap = plt.colormaps["Greys"]

    for row_idx in range(coverage.shape[0]):
        for col_idx in range(coverage.shape[1]):
            value = coverage[row_idx, col_idx]
            color = coverage_cmap(coverage_norm(value))
            ax.add_patch(plt.Rectangle((col_idx - 0.5, row_idx - 0.5), 1.0, 1.0, facecolor=color, edgecolor="white"))
            text_color = "white" if value >= 60.0 else "black"
            ax.text(col_idx, row_idx, f"{value:.1f}", ha="center", va="center", fontsize=9, color=text_color)

        ub_value = visibility_ub[row_idx]
        ub_color = ub_cmap(0.25 + 0.55 * ub_norm(ub_value))
        ax.add_patch(
            plt.Rectangle((coverage.shape[1] - 0.5, row_idx - 0.5), 1.0, 1.0, facecolor=ub_color, edgecolor="white")
        )
        text_color = "white" if ub_value >= 55.0 else "black"
        ax.text(coverage.shape[1], row_idx, f"{ub_value:.1f}", ha="center", va="center", fontsize=9, color=text_color)

    ax.set_xlim(-0.5, 5.5)
    ax.set_ylim(2.5, -0.5)
    ax.set_xticks(range(6))
    ax.set_xticklabels(["5", "10", "15", "20", "30", "Visibility UB"])
    ax.set_yticks(range(3))
    ax.set_yticklabels(["0°", "10°", "25°"])
    ax.set_xlabel("Ground-station budget")
    ax.set_ylabel("Elevation threshold")
    ax.set_title(title)
    ax.set_frame_on(False)
    ax.tick_params(length=0)

    scalar_mappable = cm.ScalarMappable(norm=coverage_norm, cmap=coverage_cmap)
    colorbar = fig.colorbar(scalar_mappable, ax=ax, fraction=0.046, pad=0.04)
    colorbar.set_label("Coverage (%)")
    fig.tight_layout()
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)


def build_portfolio_divergence_map(output_path: Path) -> None:
    metadata = json.loads(METADATA_200.read_text(encoding="utf-8"))
    candidate_sites = metadata["candidate_sites"]
    uniform_sites = set(json.loads(UNIFORM_RESULT.read_text(encoding="utf-8"))["selected_sites"])
    raster_sites = set(json.loads(RASTER_RESULT.read_text(encoding="utf-8"))["selected_sites"])

    left_only = sorted(uniform_sites - raster_sites)
    right_only = sorted(raster_sites - uniform_sites)
    shared = sorted(uniform_sites & raster_sites)
    world = _load_world_basemap()

    fig, ax = plt.subplots(figsize=(8, 5), dpi=150)
    world.plot(ax=ax, color="#dddddd", edgecolor="white", linewidth=0.35)

    categories = [
        ("Uniform-only", left_only, "#1f77b4"),
        ("Raster-only", right_only, "#d62728"),
        ("Shared", shared, "#2ca02c"),
    ]
    for label, indices, color in categories:
        if indices:
            lats = [float(candidate_sites[idx]["latitude_deg"]) for idx in indices]
            lons = [float(candidate_sites[idx]["longitude_deg"]) for idx in indices]
            ax.scatter(lons, lats, s=36, color=color, edgecolors="white", linewidths=0.4, label=label, zorder=3)

    ax.set_xlim(-180.0, 180.0)
    ax.set_ylim(-60.0, 85.0)
    ax.set_title("Selected Portfolios: Uniform vs. Raster Demand (25°, 20 stations)")
    ax.set_axis_off()
    legend_handles = [
        Line2D([0], [0], marker="o", color="w", markerfacecolor="#1f77b4", markeredgecolor="white", markersize=7, label="Uniform-only"),
        Line2D([0], [0], marker="o", color="w", markerfacecolor="#d62728", markeredgecolor="white", markersize=7, label="Raster-only"),
        Line2D([0], [0], marker="o", color="w", markerfacecolor="#2ca02c", markeredgecolor="white", markersize=7, label="Shared"),
    ]
    ax.legend(handles=legend_handles, loc="lower left", frameon=True)
    fig.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)


def build_pareto_frontier(output_path: Path) -> None:
    frame = pd.read_csv(DATA_PROCESSED / "pareto_frontier_200.csv")
    feasible = frame[frame["status"] == "Optimal"].copy()
    feasible["coverage_pct"] = 100.0 * feasible["coverage_target"]

    fig, ax = plt.subplots(figsize=(8, 5), dpi=150)
    ax.step(feasible["coverage_pct"], feasible["objective_value"], where="post", color="#1f77b4", linewidth=2.0)
    ax.scatter(feasible["coverage_pct"], feasible["objective_value"], color="#1f77b4", s=28, zorder=3)
    ax.axvline(30.0, color="#d62728", linestyle="--", linewidth=1.5)
    ymax = float(feasible["objective_value"].max()) * 1.12
    ax.text(30.8, ymax * 0.96, "Infeasible region", color="#d62728", rotation=90, va="top", ha="left")
    ax.set_xlim(0.0, 62.0)
    ax.set_ylim(0.0, ymax)
    ax.set_xlabel("Coverage target (%)")
    ax.set_ylabel("Cost")
    ax.set_title("Cost-Coverage Pareto Frontier (200-site proxy, 25° threshold)")
    ax.grid(alpha=0.25, linewidth=0.5)
    fig.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)


def build_alternative_design_comparison(output_path: Path) -> dict[str, float]:
    visibility_matrix = _load_visibility_matrix()
    candidate_coords = load_candidate_coords(METADATA_200)
    uniform_demand = np.load(DATA_PROCESSED / "demand.npy")
    raster_demand = np.load(DATA_PROCESSED / "demand_population_raster.npy")

    geo_selection = geographic_spread_selection(candidate_coords, k=20)
    greedy_selection = greedy_max_arcs_selection(visibility_matrix, k=20)

    uniform_result = json.loads(UNIFORM_RESULT.read_text(encoding="utf-8"))
    raster_result = json.loads(RASTER_RESULT.read_text(encoding="utf-8"))
    results = {
        "milp_uniform_coverage": float(uniform_result["coverage_fraction"]),
        "milp_raster_coverage": float(raster_result["coverage_fraction"]),
        "geo_spread_uniform_coverage": coverage_fraction(visibility_matrix, geo_selection, uniform_demand),
        "geo_spread_raster_coverage": coverage_fraction(visibility_matrix, geo_selection, raster_demand),
        "greedy_arcs_uniform_coverage": coverage_fraction(visibility_matrix, greedy_selection, uniform_demand),
        "greedy_arcs_raster_coverage": coverage_fraction(visibility_matrix, greedy_selection, raster_demand),
    }

    output_path.write_text(json.dumps(results, indent=2, sort_keys=True), encoding="utf-8")
    table = pd.DataFrame(
        [
            ("MILP", results["milp_uniform_coverage"], results["milp_raster_coverage"]),
            ("Geographic spread", results["geo_spread_uniform_coverage"], results["geo_spread_raster_coverage"]),
            ("Greedy max arcs", results["greedy_arcs_uniform_coverage"], results["greedy_arcs_raster_coverage"]),
        ],
        columns=["method", "uniform_coverage", "raster_coverage"],
    )
    table["uniform_coverage"] = table["uniform_coverage"].map(lambda value: f"{100.0 * value:.1f}%")
    table["raster_coverage"] = table["raster_coverage"].map(lambda value: f"{100.0 * value:.1f}%")
    print("\nAlternative design comparison")
    print(table.to_string(index=False))
    return results


def _insert_paragraph_after(paragraph: Paragraph, text: str = "") -> Paragraph:
    new_paragraph_element = OxmlElement("w:p")
    paragraph._p.addnext(new_paragraph_element)
    new_paragraph = Paragraph(new_paragraph_element, paragraph._parent)
    if text:
        new_paragraph.add_run(text)
    return new_paragraph


def _find_paragraph(document: Document, substrings: list[str]) -> Paragraph:
    for paragraph in document.paragraphs:
        text = paragraph.text.strip().lower()
        if all(substring in text for substring in substrings):
            return paragraph
    raise ValueError(f"Unable to find paragraph containing {substrings}")


def _insert_figure_after(document: Document, anchor_substrings: list[str], image_path: Path, caption: str) -> None:
    anchor = _find_paragraph(document, anchor_substrings)
    image_paragraph = _insert_paragraph_after(anchor)
    image_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    image_paragraph.add_run().add_picture(str(image_path), width=DocxInches(6.5))

    caption_paragraph = _insert_paragraph_after(image_paragraph, caption)
    caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    try:
        caption_paragraph.style = "Caption"
    except KeyError:
        pass


def _replace_alternative_design_text(document: Document, comparison: dict[str, float]) -> None:
    paragraph = _find_paragraph(document, ["alternative design comparison"])  # section header
    target = paragraph
    while target.text.strip() != "" and "alternative_design_comparison.json" not in target.text and "greedy-arc heuristics" not in target.text:
        next_element = target._p.getnext()
        if next_element is None:
            raise ValueError("Unable to locate the alternative-design placeholder paragraph")
        target = Paragraph(next_element, target._parent)

    replacement = (
        f"Under uniform demand, the MILP achieves {100.0 * comparison['milp_uniform_coverage']:.1f}% coverage "
        f"versus {100.0 * comparison['geo_spread_uniform_coverage']:.1f}% for geographic spread and "
        f"{100.0 * comparison['greedy_arcs_uniform_coverage']:.1f}% for the greedy-arc heuristic. "
        f"Under raster demand, the MILP advantage widens to {100.0 * comparison['milp_raster_coverage']:.1f}% "
        f"versus {100.0 * comparison['geo_spread_raster_coverage']:.1f}% (geographic spread) and "
        f"{100.0 * comparison['greedy_arcs_raster_coverage']:.1f}% (greedy-arc), confirming that demand-aware "
        "optimization is critical when coverage is concentrated over populated orbital passes."
    )
    target.clear()
    target.add_run(replacement)


def update_report(comparison: dict[str, float]) -> None:
    document = Document(REPORT_SOURCE)
    _insert_figure_after(
        document,
        ["table 1.", "max-coverage", "uniform demand"],
        FIGURES_DIR / "sensitivity_heatmap_uniform.png",
        "Figure 1. Heat map of uniform-demand coverage (%) across station budgets and elevation thresholds. "
        "The budget-constrained regime (0° row) and geometry-constrained regime (25° row) are visually separable.",
    )
    _insert_figure_after(
        document,
        ["table 3.", "worldpop"],
        FIGURES_DIR / "sensitivity_heatmap_raster.png",
        "Figure 2. Heat map of WorldPop raster-demand coverage (%) under the same scenario grid. "
        "Higher population concentration over visible passes raises coverage substantially relative to the uniform baseline.",
    )
    _insert_figure_after(
        document,
        ["table 4.", "uniform vs. raster demand comparison"],
        FIGURES_DIR / "portfolio_divergence_map.png",
        "Figure 3. Geographic distribution of the 20-station portfolios selected under uniform (blue) and raster (red) demand. "
        "Zero sites are shared (Jaccard = 0.00).",
    )
    _insert_figure_after(
        document,
        ["figure 3.", "geographic distribution"] if any("figure 3." in p.text.lower() for p in document.paragraphs) else ["table 4.", "uniform vs. raster demand comparison"],
        FIGURES_DIR / "pareto_frontier.png",
        "Figure 4. Cost-coverage Pareto frontier for the 200-site proxy candidate set at 25° elevation threshold. "
        "Coverage targets above 30% are infeasible under the 20-station budget constraint.",
    )
    _replace_alternative_design_text(document, comparison)
    document.save(REPORT_OUTPUT)


def _remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _fit_image_within(image_path: Path, max_width_in: float, max_height_in: float) -> tuple[float, float]:
    with Image.open(image_path) as image:
        width_px, height_px = image.size
    aspect = width_px / height_px
    width = min(max_width_in, max_height_in * aspect)
    height = width / aspect
    if height > max_height_in:
        height = max_height_in
        width = height * aspect
    return width, height


def _add_picture_centered(slide, image_path: Path, left_in: float, top_in: float, width_in: float, height_in: float) -> None:
    pic_width_in, pic_height_in = _fit_image_within(image_path, width_in, height_in)
    centered_left = left_in + (width_in - pic_width_in) / 2.0
    centered_top = top_in + (height_in - pic_height_in) / 2.0
    slide.shapes.add_picture(str(image_path), PptxInches(centered_left), PptxInches(centered_top), width=PptxInches(pic_width_in))


def _set_textbox_style(shape, *, font_size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def update_presentation() -> None:
    PRESENTATION_BACKUP.parent.mkdir(parents=True, exist_ok=True)
    if not PRESENTATION_BACKUP.exists():
        shutil.copy2(PRESENTATION_PATH, PRESENTATION_BACKUP)

    presentation = Presentation(PRESENTATION_BACKUP)

    slide6 = presentation.slides[5]
    for shape in list(slide6.shapes):
        top = shape.top / 914400
        if 1.0 <= top < 5.0:
            _remove_shape(shape)
    _add_picture_centered(slide6, FIGURES_DIR / "sensitivity_heatmap_uniform.png", 0.6, 1.0, 8.8, 3.95)

    slide8 = presentation.slides[7]
    for shape in list(slide8.shapes):
        left = shape.left / 914400
        top = shape.top / 914400
        if left >= 6.75 and 1.0 <= top < 4.9:
            _remove_shape(shape)
    _add_picture_centered(slide8, FIGURES_DIR / "portfolio_divergence_map.png", 5.0, 1.1, 4.7, 3.7)

    slide9 = presentation.slides[8]
    for shape in list(slide9.shapes):
        top = shape.top / 914400
        if 1.0 <= top < 5.0:
            _remove_shape(shape)
    _add_picture_centered(slide9, FIGURES_DIR / "sensitivity_heatmap_raster.png", 0.6, 1.0, 8.8, 3.95)

    new_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    new_slide.shapes.add_textbox(PptxInches(0.4), PptxInches(0.1), PptxInches(9.0), PptxInches(0.6)).text = "Cost-Coverage Pareto Frontier"
    title_box = new_slide.shapes[0]
    _set_textbox_style(title_box, font_size=26, bold=True)

    _add_picture_centered(new_slide, FIGURES_DIR / "pareto_frontier.png", 0.9, 0.9, 8.2, 3.6)
    text_box = new_slide.shapes.add_textbox(PptxInches(0.9), PptxInches(4.65), PptxInches(8.2), PptxInches(0.6))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = (
        "Budget is the binding constraint at permissive thresholds. Geometry is the binding constraint at 25°. "
        "The frontier surfaces this tradeoff without committing to a single target."
    )
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(16)

    footer_bar = new_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(0.0), PptxInches(5.1), PptxInches(10.0), PptxInches(0.53)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(248, 249, 251)
    footer_bar.line.color.rgb = RGBColor(248, 249, 251)
    left_footer = new_slide.shapes.add_textbox(PptxInches(0.3), PptxInches(5.14), PptxInches(5.0), PptxInches(0.3))
    left_footer.text = "IEOR 230  ·  Team Argus  ·  May 2025"
    _set_textbox_style(left_footer, font_size=10)
    right_footer = new_slide.shapes.add_textbox(PptxInches(7.5), PptxInches(5.14), PptxInches(2.3), PptxInches(0.3))
    right_footer.text = "UC Berkeley  ·  IEOR"
    _set_textbox_style(right_footer, font_size=10)

    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)
    last_slide_id = slide_ids[-1]
    slide_id_list.remove(last_slide_id)
    slide_id_list.insert(9, last_slide_id)

    presentation.save(PRESENTATION_PATH)


def main() -> None:
    FIGURES_DIR.mkdir(parents=True, exist_ok=True)
    build_sensitivity_heatmap(
        DATA_PROCESSED / "sensitivity_results.csv",
        FIGURES_DIR / "sensitivity_heatmap_uniform.png",
        "Uniform Demand: Coverage (%) by Budget and Elevation Threshold",
    )
    build_sensitivity_heatmap(
        DATA_PROCESSED / "sensitivity_results_population_raster.csv",
        FIGURES_DIR / "sensitivity_heatmap_raster.png",
        "WorldPop Raster Demand: Coverage (%) by Budget and Elevation Threshold",
    )
    build_portfolio_divergence_map(FIGURES_DIR / "portfolio_divergence_map.png")
    build_pareto_frontier(FIGURES_DIR / "pareto_frontier.png")
    comparison = build_alternative_design_comparison(DATA_PROCESSED / "alternative_design_comparison.json")
    update_report(comparison)
    update_presentation()


if __name__ == "__main__":
    main()
